"""
Build the results slide deck (deliverable) from the figures the pipeline saves.

Reproducible: re-run after `--all` to rebuild results/slides/GI_nanoplastic.pptx
from whatever is currently in results/full/figures/. Missing figures are skipped with
a note rather than crashing, so a partial run still yields a deck.

    python -m src.make_slides          # or:  .\run.ps1 slides
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

from . import config as cfg

SLIDES_DIR = cfg.PROJECT_ROOT / "results" / "slides"
SLIDES_DIR.mkdir(parents=True, exist_ok=True)

# 16:9 widescreen (the modern default). Layout math below derives from these
# so figure placement stays correct if the dimensions ever change.
SLIDE_W, SLIDE_H = 13.333, 7.5
MARGIN = 0.4   # left/right page margin (inches)

# The deck is a deliverable built from the FULL run's figures. Resolve to the
# full-run location explicitly (independent of any smoke run's cfg state).
FIG = cfg.RESULTS_DIR / "full" / "figures"


def _stretch(ph):
    """Widen a placeholder to the usable width WITHOUT disturbing its vertical
    position. python-pptx's built-in layouts size placeholders for a 4:3 slide
    (~9in wide); on our 16:9 slide that leaves text hugging the left with a dead
    band on the right. Capture the inherited top/height first: setting left/width
    materializes the shape's xfrm, and any position attribute left unset then
    defaults to 0 -- which previously yanked every text box to top=0, height=0."""
    top, height = ph.top, ph.height        # inherited from the layout; read before mutating
    ph.left = Inches(MARGIN)
    ph.width = Inches(SLIDE_W - 2 * MARGIN)
    ph.top = top
    ph.height = height


def _notes(slide, text):
    """Attach speaker notes (visible in PowerPoint's Presenter View, not on the
    slide itself). Used to carry the narration + 'what you're seeing' per slide."""
    if text:
        slide.notes_slide.notes_text_frame.text = text.strip()


def _title_slide(prs, title, subtitle, notes=""):
    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = title
    s.placeholders[1].text = subtitle
    for ph in s.placeholders:
        _stretch(ph)
    _notes(s, notes)


def _bullets_slide(prs, title, bullets, notes=""):
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = title
    _stretch(s.shapes.title)
    body = s.placeholders[1]
    _stretch(body)
    tf = body.text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.font.size = Pt(18)
    _notes(s, notes)


def _place_fit(slide, path, box_left, box_top, box_w, box_h):
    """Add a picture scaled to FIT inside (box_w x box_h) without distortion or
    overflow -- letterbox by whichever dimension binds, then centre it in the box.
    Sizing by height alone (the old behaviour) let wide figures (violin strips
    are 3.75:1, additivity 5:1) blow past the slide edge and overlap."""
    from PIL import Image
    px_w, px_h = Image.open(str(path)).size
    aspect = px_w / px_h            # width / height of the actual image
    box_aspect = box_w / box_h
    if aspect >= box_aspect:
        w, h = box_w, box_w / aspect   # width-bound (wide image)
    else:
        w, h = box_h * aspect, box_h   # height-bound (tall/square image)
    left = box_left + (box_w - w) / 2  # centre within the box
    top = box_top + (box_h - h) / 2
    slide.shapes.add_picture(str(path), Inches(left), Inches(top),
                             width=Inches(w), height=Inches(h))


def _figure_slide(prs, title, fig_names, caption="", notes=""):
    """One slide, title + up to two figures side by side (skips any missing).
    Positions derive from SLIDE_W / MARGIN so the layout tracks the slide size."""
    paths = [FIG / f for f in fig_names if (FIG / f).exists()]
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = title
    _stretch(s.shapes.title)
    usable_w = SLIDE_W - 2 * MARGIN
    if not paths:
        tb = s.shapes.add_textbox(Inches(MARGIN), Inches(2.5), Inches(usable_w), Inches(1)).text_frame
        tb.text = "(figure not generated yet -- run the pipeline)"
        _notes(s, notes)
        return
    top, box_h = 1.5, 5.0            # vertical band for figures (below title, above caption)
    if len(paths) == 1:
        _place_fit(s, paths[0], MARGIN, top, usable_w, box_h)   # one figure, full usable width
    else:
        gap = 0.3
        half = (usable_w - gap) / 2
        _place_fit(s, paths[0], MARGIN, top, half, box_h)              # left half
        _place_fit(s, paths[1], MARGIN + half + gap, top, half, box_h) # right half
    if caption:
        tb = s.shapes.add_textbox(Inches(MARGIN), Inches(6.9), Inches(usable_w), Inches(0.5)).text_frame
        tb.text = caption
        tb.paragraphs[0].font.size = Pt(12)
    _notes(s, notes)


def build():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    _title_slide(prs, "Single-Cell Immune Response to Nanoplastic Particles",
                 "scRNA-seq of human PBMCs exposed to 40 nm / 200 nm / mixed polystyrene "
                 "nanoparticles vs control  -  Genomics Informatics, ETF",
                 notes=(
                     "Hi, my name is Vuk Djordjevic. I'm a master's student at the School of "
                     "Electrical Engineering in Belgrade, and this is my project for the Genome "
                     "Informatics course, 2026. The project is a single-cell analysis of how human "
                     "immune cells respond to nanoplastic particles of different sizes. Nanoplastics "
                     "are tiny plastic particles, and they've now been found in human blood, in "
                     "direct contact with our immune cells. The question I looked at is whether "
                     "particle size matters: do small and large nanoplastics cause different immune "
                     "responses, and does a mixture of both do something neither size does on its "
                     "own? To answer this I used single-cell RNA sequencing, which lets me read "
                     "which genes are active in each individual cell."))
    _bullets_slide(prs, "Question & dataset", [
        "Do small vs large nanoplastics provoke different immune responses?",
        "Does the 40+200 nm mixture do something neither size does alone?",
        "4 samples, one donor: PSNP_40nm, PSNP_200nm, PSNP_mixture, control",
        "~34,000 PBMCs, AnnData/.h5ad (Zenodo 10.5281/zenodo.15866724)",
    ], notes=(
        "The data is four samples from a single donor. Peripheral blood immune cells were exposed "
        "to 40-nanometre particles, to 200-nanometre particles, to a mixture of both, and one "
        "sample was left unexposed as a control. That's about 34,000 cells in total. One thing I "
        "want to be upfront about: there is only one donor and one sample per condition, so there "
        "are no biological replicates. That limitation shaped how I did the statistics, and I flag "
        "it throughout the project."))
    _bullets_slide(prs, "Pipeline (6 stages, one reproducible driver)", [
        "1. QC & preprocessing  -  thresholds justified on real percentiles",
        "2. Integration  -  Harmony batch correction, UMAP, Leiden",
        "3. Annotation  -  celltypist, cross-checked vs Azimuth & CoDi (~93%)",
        "4. Composition  -  proportions + log2 fold-change vs control",
        "5. Differential expression  -  per-lineage Wilcoxon + pathway enrichment",
        "6. Size-specific effects  -  unique / shared / mixture-emergent genes",
    ], notes=(
        "The whole analysis runs through one pipeline of six stages: quality control, integration, "
        "cell-type annotation, composition analysis, differential expression, and size-specific "
        "effects. A single command reproduces all of it from the raw data, and each stage has "
        "automated tests."))
    # QC split across two slides -- the violin strip (3.75:1) and the scatter
    # are both too wide to share one slide without shrinking to unreadable.
    _figure_slide(prs, "QC & preprocessing (1/2): per-sample QC metrics", ["01_qc_violin_after.png"],
                  "Violins of the 3 QC metrics after filtering. Thresholds justified on the real distribution "
                  "(max_genes ~p99, max_mito between p95-p99).",
                  notes=(
                      "The first stage is quality control, where I remove low-quality cells before trusting "
                      "anything. These three plots are violins, which are just histograms turned on their side "
                      "and mirrored, so they're fatter where more cells sit. Along the bottom of each are my "
                      "four samples. The left plot is how many genes each cell expressed, the middle is the "
                      "total reads per cell, and the right is the percent of reads coming from mitochondria, "
                      "which is a sign of a dying cell. This is after cleaning, and you can see the shapes look "
                      "healthy and similar across all four samples, so no sample is junk. Instead of guessing "
                      "the cutoffs I set them from the real distribution: I drop cells with too many genes, "
                      "which are usually two cells stuck together, and cells with too much mitochondria, which "
                      "are dying."))
    _figure_slide(prs, "QC & preprocessing (2/2): counts vs %mito", ["01_qc_scatter_counts_mito.png"],
                  "total_counts vs %mito (pre-filter); threshold lines drawn. Drops likely doublets and dying cells.",
                  notes=(
                      "This is the same cleaning step, drawn a different way. Every dot here is one cell. "
                      "Left-to-right is how many reads the cell has, and bottom-to-top is how mitochondrial it "
                      "is. The red dashed line is my cutoff at 15 percent. Almost all the cells form a dense "
                      "cloud hugging the bottom, with low mitochondria, which is exactly what healthy cells "
                      "look like. The handful of dots floating up above the red line are dying or broken cells, "
                      "and those get thrown out."))
    _figure_slide(prs, "Integration removes the batch effect", ["02_umap_pre_harmony_by_sample.png", "02_umap_by_sample.png"],
                  "Before (left) vs after Harmony (right), coloured by sample: samples mix after correction.",
                  notes=(
                      "The second stage is integration. Both of these pictures are a UMAP. The way to picture "
                      "it: each cell has thousands of genes, and UMAP squashes all of that down to a single dot "
                      "on a map, so that cells with similar gene activity land close together. The axes are "
                      "just map coordinates with no units; only closeness means anything. Each dot is a cell, "
                      "coloured by which sample it came from. What I want here is for the colours to be well "
                      "mixed. If the cells clumped up by colour, it would mean the sample a cell came from "
                      "matters more than its actual biology, which is a technical artefact I'd have to remove. "
                      "On both panels the colours are already nicely blended, so the samples line up and the "
                      "cells group by cell type rather than by batch. One honest caveat: because each condition "
                      "is a single sample, the batch and the treatment effects are tied together, so for the "
                      "gene-level results later I work on the uncorrected data, where this can't bias "
                      "anything."))
    _figure_slide(prs, "Cell-type annotation", ["03_umap_lineage.png", "03_marker_dotplot.png"],
                  "celltypist lineages; agreement with Azimuth 92.7% and CoDi 93.1% (independent references).",
                  notes=(
                      "The third stage is labelling the cell types. On the left is the same cell-map, now "
                      "coloured by what kind of immune cell each dot is. You can literally see the "
                      "neighbourhoods: the T cells are the big purple mass on the right, the monocytes are the "
                      "orange island top-left, the B cells are blue at the bottom, and the NK cells are green. "
                      "The fact that they form clean separate islands means the cell types are real and "
                      "distinct. On the right is a dot-plot that proves the labels. The rows are the cell "
                      "clusters and the columns are famous marker genes, like CD3D for T cells or CD14 for "
                      "monocytes. A big dark dot means that gene is strongly switched on in that cluster, and "
                      "you can see the dark dots line up on the diagonal, so each cluster lights up its own "
                      "known marker. I assigned the types with a tool called celltypist, then cross-checked "
                      "against two independent methods already in the data, and my labels agreed with both at "
                      "about 93 percent."))
    _figure_slide(prs, "Composition shifts vs control", ["04_composition_stacked.png", "04_composition_grouped.png"],
                  "Cell-type proportions per sample; PSNP_200nm shows the largest compositional shift.",
                  notes=(
                      "The fourth stage asks whether exposure changed how many of each cell type there are. "
                      "The bars show what fraction of a sample is each cell type, and the four coloured bars in "
                      "each group are the four samples. The T cells dominate every sample, those are the tall "
                      "bars on the right at around 80 percent, which is normal for blood. The thing to notice "
                      "is the Monocyte group: the green 200-nanometre bar is noticeably taller than the others, "
                      "roughly two and a half times the control's monocyte fraction, while everything else "
                      "barely moves. So 200-nanometre exposure specifically bumps up the monocytes. Since there "
                      "are no replicates, I report this as a descriptive shift and treat any statistical test "
                      "as exploratory."))
    _figure_slide(prs, "Differential expression: dose-response", ["07_dose_response.png", "06_size_categories.png"],
                  "Monocytes respond most; 200 nm drives more unique genes than 40 nm; mixture-emergent in monocytes.",
                  notes=(
                      "The fifth stage asks which genes actually change on exposure, within each cell type. "
                      "There are two bar charts, and in both the height of a bar is the number of genes that "
                      "significantly changed, so taller means a bigger reaction. The left chart groups by "
                      "condition and colours by cell type. The story jumps right out: the orange Monocyte bars "
                      "tower over everything, because monocytes are the cells that eat and clear foreign "
                      "particles so they react the hardest, and the 200-nanometre bar is taller than the "
                      "40-nanometre one, meaning bigger particles cause a bigger response. The right chart is "
                      "the sixth stage, and it zooms into where those genes overlap. For each cell type there "
                      "are four bars: genes that react only to 40 nanometres, only to 200, to both, and the "
                      "key one in red, genes that fire only in the mixture and in neither size on its own. Look "
                      "at the Monocyte group: that red bar is about 180 genes, a whole response that only "
                      "appears when both sizes are combined. So the mixture does something the individual sizes "
                      "do not."))
    _bullets_slide(prs, "Key biological finding", [
        "Monocytes show the strongest transcriptional response to nanoplastic.",
        "200 nm particles drive MORE lineage-unique genes than 40 nm.",
        "The 40+200 nm mixture produces an EMERGENT monocyte response -",
        "  genes significant only in the mixture, in neither single size.",
        "In lymphocytes the mixture is weaker than either size (sub-additive).",
    ], notes=(
        "So let me pull the main findings together. First, monocytes show the strongest response "
        "to nanoplastic of any cell type. Second, the 200-nanometre particles drive more "
        "cell-type-unique genes than the 40-nanometre ones. And third, the biggest one: the 40 "
        "plus 200-nanometre mixture produces an emergent monocyte response, about 180 genes that "
        "are significant only in the mixture and in neither single size. Interestingly, in the "
        "lymphocytes the opposite happens, the mixture is actually weaker than either size alone. "
        "Either way, the combination of the two sizes does something the individual sizes do "
        "not."))
    _bullets_slide(prs, "What the mixture does: emergent inflammation", [
        "The 180 mixture-only monocyte genes are dominated by INFLAMMATION.",
        "Pathway enrichment (p ~ 1e-21): interleukin, cytokine & TNF signalling,",
        "  and 'response to lipopolysaccharide' - i.e. a bacterial-attack-like program.",
        "Top mixture-only genes up: RIPK2, TRAF1, PIK3CB (innate/TNF signalling).",
        "Interpretation: two sub-threshold exposures together push monocytes over",
        "  an inflammatory activation threshold - an emergent, non-additive effect.",
        "Real exposure is always to mixtures -> single-size studies may under-estimate risk.",
    ], notes=(
        "So what is that new thing the mixture does? I ran pathway enrichment on those 180 "
        "mixture-only monocyte genes, and the result points clearly at inflammation. The top "
        "pathways, with very strong statistical support, are interleukin, cytokine, and TNF "
        "signalling, along with the response to lipopolysaccharide. That last one is the program a "
        "monocyte runs during a bacterial infection, except here it's set off by plastic. Two of "
        "the strongest genes that go up in the mixture are RIPK2 and TRAF1, which are core "
        "innate-immune and TNF-signalling genes. My reading is that 40 and 200 nanometres each "
        "cause small changes that stay below a threshold, but together they push the monocyte past "
        "that threshold into an inflammatory state. And since real-world exposure is always to a "
        "mix of particle sizes, testing single sizes alone could underestimate the risk."))
    # Bonus split across two slides -- the additivity plot (5:1) is far too wide
    # to sit beside the module-scores heatmap on one slide.
    _bullets_slide(prs, "Additional analyses (5 implemented)", [
        "1. Module scoring - stress & inflammation gene programs per sample (shown next).",
        "2. Mixture additivity - is the mixture = 40nm + 200nm? (shown next).",
        "3. Clustering robustness - clusters stable across Leiden resolutions (ARI).",
        "4. Ligand-receptor - shift in cell-to-cell communication on exposure.",
        "5. Dose-response - disruption magnitude vs particle size.",
        "All five reinforce the same story: size-dependent, non-additive, monocyte-centred.",
    ], notes=(
        "For the additional-insights part of the project I implemented five extra analyses, which "
        "this slide lists, and the next slides show each one. They are module scoring, a mixture "
        "additivity test, a clustering robustness check, a ligand-receptor communication analysis, "
        "and dose-response. I'll walk through them in the next few slides. All five point the same "
        "way: a size-dependent, non-additive response centred on the monocytes."))
    _figure_slide(prs, "Additional analyses: stress/inflammation module scores", ["07_module_scores.png"],
                  "Analysis 1 of 5. Per-sample mean module scores for stress and inflammation gene programs.",
                  notes=(
                      "The first additional analysis is module scoring. This is a colour grid where each row "
                      "is a sample and each column is a biological program, which is just a named group of "
                      "genes: oxidative stress, NF-kB inflammation, interferon, and heat shock. Red means that "
                      "program is turned up in that sample, and white or blue means it isn't. The whole left "
                      "column, oxidative stress, goes deep red in all three exposed samples but not in the "
                      "control, so the plastic reliably stresses the cells. And the inflammation column turns "
                      "pink specifically under 200 nanometres. So the size effect I keep pointing at shows up "
                      "here too, in independent gene programs."))
    _figure_slide(prs, "Additional analyses: mixture additivity", ["07_mixture_additivity.png"],
                  "Analysis 2 of 5. Observed mixture response vs the 40+200 nm additive expectation, per lineage.",
                  notes=(
                      "The second analysis is a direct test of whether the mixture is just 40 plus 200 added "
                      "together. There's one scatter per cell type, and each dot is a gene. Left-to-right is "
                      "what I'd predict if the mixture were simply the two sizes summed, and bottom-to-top is "
                      "what the mixture actually did. The red diagonal line is where the prediction would be "
                      "perfect. If every gene sat on that line, the mixture would be boringly additive. But "
                      "look at the Monocyte panel: there are clear blobs of genes sitting off the line. Those "
                      "are genes the mixture switched on or off that you'd never predict from the single sizes, "
                      "and that off-the-line behaviour is exactly the emergent, non-additive effect."))
    _figure_slide(prs, "Additional analyses: robustness & communication",
                  ["07_clustering_robustness.png", "07_ligand_receptor.png"],
                  "Analyses 3 & 4 of 5. Left: clustering robustness (ARI across Leiden resolutions). "
                  "Right: ligand-receptor communication shift on exposure. (Analysis 5, dose-response, is on the DE slide.)",
                  notes=(
                      "These are the third and fourth analyses. The plot on the left is a sanity check on my "
                      "clustering. The bottom axis is how finely I chose to cut the cells into groups. The blue "
                      "line just rises because you naturally get more clusters if you cut finer. The one that "
                      "matters is the green dashed line, which is how much the grouping still agrees with my "
                      "default setting, where 1.0 would be identical. It stays high near the setting I chose, "
                      "so my cell groups aren't a fluke of one knob. The plot on the right is cell-to-cell "
                      "messaging. Each row is a signalling pair, a messenger and its receptor, and a bar goes "
                      "right if that message got louder after exposure. The standout is IL1B, a classic "
                      "inflammation alarm signal, shooting far right in all the exposed samples and furthest "
                      "under 200 nanometres. So the cells aren't just changing internally, they're shouting "
                      "inflammation at each other. The fifth analysis, dose-response, I already showed on the "
                      "differential-expression slide earlier."))
    _bullets_slide(prs, "Limitations & reproducibility", [
        "One donor, one sample per condition -> NO biological replicates.",
        "DE is cell-level Wilcoxon (not pseudobulk); p-values exploratory, pseudoreplication caveat.",
        "No external gene-level ground truth; DE validated internally + biologically.",
        "Everything reproduces: `python run_pipeline.py --all`; 37 offline intent tests.",
    ], notes=(
        "Finally, to be clear about the limits. One donor and no replicates mean the differential "
        "expression is done per cell and is exploratory, with a pseudoreplication caveat. There's "
        "no external gene-level reference, so I validated the results internally and biologically. "
        "But everything reproduces from the repository with a single command and a full test "
        "suite. Thanks for watching."))

    out = SLIDES_DIR / "GI_nanoplastic.pptx"
    prs.save(out)
    print(f"saved slide deck -> {out}  ({len(list(prs.slides))} slides)")
    return out


if __name__ == "__main__":
    build()
